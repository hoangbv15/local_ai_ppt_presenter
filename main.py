#!/usr/bin/env python
# -*- coding: utf-8 -*-

import os
from pathlib import Path
import tempfile
import argparse
from subprocess import call

from pdf2image import convert_from_path
from pptx import Presentation
from ttsgen import TTSGen
from engines.xtts2_engine import XTTS2Engine
from engines.gtts_engine import GTTSEngine

__author__ = ['hoangbv15']

## Sometimes ffmpeg is avconv
FFMPEG_NAME = 'ffmpeg'
#FFMPEG_NAME = 'avconv'


def ppt_presenter(pptx_path, pdf_path, output_path, temp_dir, engineName, fast, saveclips, pagenos, saveaudio):
    if fast:
        tts = TTSGen(GTTSEngine())
    elif engineName:
        engine = globals()[engineName]
    else:
        tts = TTSGen(XTTS2Engine())

    with tempfile.TemporaryDirectory(dir=temp_dir) as temp_path:
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)

        tts.enable(True)
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            if pagenos and i not in pagenos:
                continue

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                audio_path = os.path.join(temp_path, 'frame_{}.wav'.format(i+1))
                tts.generate(text=notes,
                             output_file=audio_path)
                if saveaudio:
                    continue
                image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i+1))
                image.save(image_path)

                ffmpeg_call(image_path, audio_path, temp_path, i+1)

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i+1)) \
                      for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)

        if saveclips or saveaudio:
            output_path = output_path.replace('.mp4', '-clips')
            print("saveclips option is set")
            if not os.path.exists(output_path):
                os.makedirs(output_path)

            src_path = Path(temp_path)
            dest_path = Path(output_path)
            glob = src_path.glob('*.mp4')
            if saveaudio:
                glob = src_path.glob('*.wav')
            for each_file in glob:
                print("Moving % s to % s" % (each_file.name, output_path))
                each_file.rename(dest_path.joinpath(each_file.name))

def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
    call([FFMPEG_NAME, '-loop', '1', '-y', '-i', image_path, '-i', audio_path, '-shortest', '-fflags', '+shortest', '-max_interleave_delta', '200M',
      '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'aac',
      '-b:a', '192k', '-vf', 'scale=-1:1080', out_path_mp4])
    call([FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])


def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])


def main():
    parser = argparse.ArgumentParser(description='Local AI PPT Presenter help.')
    parser.add_argument('--pptx', help='input pptx path')
    parser.add_argument('--pdf', help='input pdf path')
    parser.add_argument('-o', '--output', help='output path')
    parser.add_argument('-t', '--tempdir', help='path to store temporary files needed to generate the output. A ramdisk is recommended. Leave none to use python tempfile defaults.')
    parser.add_argument('-e', '--engine', help='the name of the text to speech engine to use')
    parser.add_argument('-f', '--fast', help='use the text to speech engine of the OS for fast execution at the expense of quality', action='store_true')
    parser.add_argument('-sc', '--saveclips', help='save the clips for each page instead of deleting them', action='store_true')
    parser.add_argument('-sa', '--saveaudio', help='only generate voice audios', action='store_true')
    parser.add_argument('-p', '--pageno', help='only regenerate the given page number')
    args = parser.parse_args()

    pagenos = []
    if args.pageno:
        for page in args.pageno.split(','):
            r = page.split('-')
            if len(r) > 1:
                for i in range(int(r[0]) - 1, int(r[1])):
                    pagenos.append(i)
                continue
            pagenos.append(int(page) - 1)
        args.saveclips = True
        print('Page list: % s' % pagenos)

    ppt_presenter(args.pptx, args.pdf, args.output, args.tempdir, 
                args.engine, args.fast, args.saveclips, pagenos, args.saveaudio)


if __name__ == '__main__':
    main()